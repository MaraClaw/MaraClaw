"""Extract text from common office file formats.

Supports: PDF, DOCX, XLSX, PPTX
Saves extracted text as a companion .md file alongside the original.
"""

import io
from collections.abc import Iterable, Sequence
from pathlib import Path
from typing import Protocol, TypeIs

from app.core.logging import logger

# File extensions that need text extraction
EXTRACTABLE_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Text extensions that don't need extraction
TEXT_EXTS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".js",
    ".ts",
    ".py",
    ".html",
    ".css",
    ".sh",
    ".log",
    ".env",
}


class _PdfContext(Protocol):
    def __enter__(self) -> object: ...

    def __exit__(self, *args: object) -> object: ...


class _PdfPlumberModule(Protocol):
    def open(self, stream: object) -> _PdfContext: ...


class _DocxModule(Protocol):
    def Document(self, stream: object) -> object: ...


class _Workbook(Protocol):
    sheetnames: object

    def __getitem__(self, name: str) -> object: ...

    def close(self) -> object: ...


class _OpenpyxlModule(Protocol):
    def load_workbook(self, stream: object, read_only: bool = False, data_only: bool = False) -> _Workbook: ...


class _PptxModule(Protocol):
    def Presentation(self, stream: object) -> object: ...


class _Nullary(Protocol):
    def __call__(self, *args: object, **kwargs: object) -> object: ...


def _is_nullary(value: object) -> TypeIs[_Nullary]:
    return callable(value)


def _call_sdk(value: object, *args: object, **kwargs: object) -> object:
    if not _is_nullary(value):
        return None
    return value(*args, **kwargs)


def _is_pdfplumber_module(value: object) -> TypeIs[_PdfPlumberModule]:
    return callable(getattr(value, "open", None))


def _is_docx_module(value: object) -> TypeIs[_DocxModule]:
    return callable(getattr(value, "Document", None))


def _is_openpyxl_module(value: object) -> TypeIs[_OpenpyxlModule]:
    return callable(getattr(value, "load_workbook", None))


def _is_pptx_module(value: object) -> TypeIs[_PptxModule]:
    return callable(getattr(value, "Presentation", None))


def _object_sequence(value: object) -> list[object]:
    if isinstance(value, Sequence) and not isinstance(value, (str, bytes)):
        return list[object](value)
    return []


def _object_iterable(value: object) -> list[object]:
    if isinstance(value, (str, bytes)):
        return []
    if isinstance(value, Sequence):
        return list[object](value)
    if isinstance(value, Iterable):
        return list[object](value)
    return []


def _attr_text(value: object, name: str = "text") -> str:
    text: object = getattr(value, name, "")
    return text.strip() if isinstance(text, str) else str(text or "").strip()


def _clean_cell(value: object) -> str:
    text = str(value or "").strip()
    return text.replace("\n", "<br>").replace("|", "\\|")


def _markdown_table(rows: Sequence[Sequence[object]]) -> str:
    cleaned = [[_clean_cell(cell) for cell in row] for row in rows]
    cleaned = [row for row in cleaned if any(cell for cell in row)]
    if not cleaned:
        return ""

    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in cleaned]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def needs_extraction(filename: str) -> bool:
    """Check if a file needs text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in EXTRACTABLE_EXTS


def extract_text(file_bytes: bytes, filename: str) -> str | None:
    """Extract text from a binary file.

    Returns extracted text string, or None if extraction fails.
    """
    ext = Path(filename).suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        if ext == ".docx":
            return _extract_docx(file_bytes)
        if ext == ".xlsx":
            return _extract_xlsx(file_bytes)
        if ext == ".pptx":
            return _extract_pptx(file_bytes)
    except Exception as e:
        logger.error(f"[TextExtractor] Failed to extract from {filename}: {e}")
        return None

    return None


def save_extracted_text(save_path: Path, file_bytes: bytes, filename: str) -> Path | None:
    """Extract text and save as a companion .md file.

    For example: report.pdf -> report.md
    Returns the path to the markdown file, or None if extraction failed.
    """
    text = extract_text(file_bytes, filename)
    if not text or not text.strip():
        return None

    md_path = save_path.parent / f"{save_path.stem}.md"
    _ = md_path.write_text(text, encoding="utf-8")
    logger.info(f"[TextExtractor] Extracted {len(text)} chars from {filename} -> {md_path.name}")
    return md_path


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber

    pdfplumber_mod: object = pdfplumber
    if not _is_pdfplumber_module(pdfplumber_mod):
        raise TypeError("pdfplumber.open is unavailable")

    pages: list[str] = []
    with pdfplumber_mod.open(io.BytesIO(data)) as pdf:
        pdf_obj: object = pdf
        for i, page in enumerate(_object_sequence(getattr(pdf_obj, "pages", []))):
            page_parts: list[str] = []
            text = _call_sdk(getattr(page, "extract_text", None))
            page_text = text.strip() if isinstance(text, str) else ""
            if page_text:
                page_parts.append(page_text)

            tables = _call_sdk(getattr(page, "extract_tables", None))
            for table in _object_sequence(tables):
                table_rows = [_object_sequence(row) for row in _object_sequence(table)]
                table_md = _markdown_table(table_rows)
                if table_md:
                    page_parts.append(table_md)

            if page_parts:
                pages.append(f"## Page {i + 1}\n\n" + "\n\n".join(page_parts))

    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    import docx

    docx_mod: object = docx
    if not _is_docx_module(docx_mod):
        raise TypeError("docx.Document is unavailable")

    doc: object = docx_mod.Document(io.BytesIO(data))
    parts: list[str] = []

    for para in _object_sequence(getattr(doc, "paragraphs", [])):
        text = _attr_text(para)
        if text:
            style: object = getattr(para, "style", None)
            style_name_raw: object = getattr(style, "name", "")
            style_name = style_name_raw if isinstance(style_name_raw, str) else ""
            if style_name.startswith("Heading"):
                level_raw = style_name.replace("Heading", "").strip()
                try:
                    level = int(level_raw)
                except ValueError:
                    level = 1
                parts.append(f"{'#' * level} {text}")
            elif "List Bullet" in style_name:
                parts.append(f"- {text}")
            elif "List Number" in style_name:
                parts.append(f"1. {text}")
            else:
                parts.append(text)

    for table in _object_sequence(getattr(doc, "tables", [])):
        rows = [
            [_attr_text(cell) for cell in _object_sequence(getattr(row, "cells", []))]
            for row in _object_sequence(getattr(table, "rows", []))
        ]
        table_md = _markdown_table(rows)
        if table_md:
            parts.append("## Table\n\n" + table_md)

    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    import openpyxl

    openpyxl_mod: object = openpyxl
    if not _is_openpyxl_module(openpyxl_mod):
        raise TypeError("openpyxl.load_workbook is unavailable")

    wb = openpyxl_mod.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []

    sheet_names = [name for name in _object_sequence(wb.sheetnames) if isinstance(name, str)]
    for sheet in sheet_names:
        ws: object = wb[sheet]
        raw_rows = _call_sdk(getattr(ws, "iter_rows", None), values_only=True)
        rows: list[list[object]] = []
        for row in _object_iterable(raw_rows):
            if not isinstance(row, Sequence) or isinstance(row, (str, bytes)):
                continue
            cells = [cell if cell is not None else "" for cell in list[object](row)]
            if any(str(cell).strip() for cell in cells):
                rows.append(cells)

        table_md = _markdown_table(rows)
        if table_md:
            parts.append(f"## Worksheet: {sheet}\n\n" + table_md)

    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    import pptx

    pptx_mod: object = pptx
    if not _is_pptx_module(pptx_mod):
        raise TypeError("pptx.Presentation is unavailable")

    prs: object = pptx_mod.Presentation(io.BytesIO(data))
    parts: list[str] = []

    for i, slide in enumerate(_object_sequence(getattr(prs, "slides", []))):
        texts: list[str] = []
        tables: list[str] = []
        for shape in _object_sequence(getattr(slide, "shapes", [])):
            text_frame = getattr(shape, "text_frame", None) if getattr(shape, "has_text_frame", False) else None
            if text_frame is not None:
                for para in _object_sequence(getattr(text_frame, "paragraphs", [])):
                    text = _attr_text(para)
                    if text:
                        texts.append(text)
            table = getattr(shape, "table", None) if getattr(shape, "has_table", False) else None
            if table is not None:
                rows = [
                    [_attr_text(cell) for cell in _object_sequence(getattr(row, "cells", []))]
                    for row in _object_sequence(getattr(table, "rows", []))
                ]
                table_md = _markdown_table(rows)
                if table_md:
                    tables.append(table_md)

        slide_parts: list[str] = []
        if texts:
            slide_parts.append("\n\n".join(texts))
        slide_parts.extend(tables)
        if slide_parts:
            parts.append(f"## Slide {i + 1}\n\n" + "\n\n".join(slide_parts))

    return "\n\n".join(parts)
